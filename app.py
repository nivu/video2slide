import os
import time
import cv2
import imutils
import shutil
import time
import pytesseract 
from pptx import Presentation
from pptx.util import Inches
import youtube_dl
from tkinter import messagebox, filedialog
import streamlit as st

st.set_page_config(page_title ="QA Gen",page_icon="☁️")

html_temp = """
        <div style="background-color:skyblue;padding:10px">
        <h1 style="color:white;text-align:center;">VIDEO TO SLIDE CONVERTER </h1>
        </div>
        """

OUTPUT_SLIDES_DIR = f"./output"
framerate = 3                   # no.of frames per second that needs to be processed, fewer the count faster the speed
warmup = framerate              # initial number of frames to be skipped
fgbshistory = framerate * 15    # no.of frames in background object
varthreshold = 16               # Threshold on the squared Mahalanobis distance between the pixel and the model to decide whether a pixel is well described by the background model.
detectshadows = False           # If true, the algorithm will detect shadows and mark them.
minipercent = 0.1               # min % of diff between foreground and background to detect if motion has stopped
maxpercent = 3                  # max % of diff between foreground and background to detect if frame is still in motion

def get_frames(video_path):

    vs = cv2.VideoCapture(video_path)
    framerate=vs.get(cv2.CAP_PROP_FPS)
    if not vs.isOpened():
        raise Exception(f'unable to open file {video_path}')

    total_frames = vs.get(cv2.CAP_PROP_FRAME_COUNT)
    frame_time = 0
    frame_count = 0
    print("total_frames: ", total_frames)
    print("framerate", framerate)

    while True:

        vs.set(cv2.CAP_PROP_POS_MSEC, frame_time * 1000)    # move frame to a timestamp
        frame_time += 1/framerate

        (_, frame) = vs.read()

        if frame is None:
            break

        frame_count += 1
        yield frame_count, frame_time, frame

    vs.release()

pytesseract.pytesseract.tesseract_cmd = 'C:\Program Files\Tesseract-OCR/tesseract.exe'

def detect_unique_screenshots(video_path, output_folder_screenshot_path):

    fgbg = cv2.createBackgroundSubtractorMOG2(history=warmup, varThreshold=varthreshold,detectShadows=detectshadows)

    captured = False
    start_time = time.time()
    (W, H) = (None, None)

    screenshoots_count = 0
    for frame_count, frame_time, frame in get_frames(video_path):
        orig = frame.copy() # clone the original frame (so we can save it later), 
        frame = imutils.resize(frame, width=600) # resize the frame
        mask = fgbg.apply(frame) # apply the background subtractor

        if W is None or H is None:
            (H, W) = mask.shape[:2]

        p_diff = (cv2.countNonZero(mask) / float(W * H)) * 100

        if p_diff < minipercent and not captured and frame_count > warmup:
            captured = True
            filename = f"{screenshoots_count:03}_{round(frame_time/60, 2)}.png"

            path = os.path.join(output_folder_screenshot_path, filename)
            print("saving {}".format(path))
            cv2.imwrite(path, orig)
            img1 = cv2.imread(path)
            img2 = cv2.cvtColor(img1, cv2.COLOR_BGR2RGB)
            convert_screenshots_to_pptx(img2,path)
            screenshoots_count += 1

        elif captured and p_diff >= maxpercent:
            captured = False
    print(f'{screenshoots_count} screenshots Captured!')
    print(f'Time taken {time.time()-start_time}s')
    return 

def initialize_output_folder(video_path):
    '''Clean the output folder if already exists'''
    output_folder_screenshot_path = f"{OUTPUT_SLIDES_DIR}/{video_path.rsplit('/')[-1].split('.')[0]}"

    if os.path.exists(output_folder_screenshot_path):
        shutil.rmtree(output_folder_screenshot_path)

    os.makedirs(output_folder_screenshot_path, exist_ok=True)
    print('initialized output folder', output_folder_screenshot_path)
    return output_folder_screenshot_path

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def convert_screenshots_to_pptx(img2,path):
    mytext=""
    i=0
    k=[]
    texts =  pytesseract.image_to_data(img2) 
    img_path: str = path
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    top = Inches(1)
    left1 = Inches(12)
    height1 = width1 = Inches(3.5)
    for b,o in enumerate(texts.splitlines()):
        g=o.split()
        k.append(g)
    g.append("")
    g[6]=0
    k.append(g)
    m=10000
    lst1=[1]
    lst2=[1]
    for j in k:
        if len(j)==12 and i>=1:
            x,y,w,h = int(j[6]),int(j[7]),int(j[8]),int(j[9])
            if x<m and i>=2:
                txBox = slide.shapes.add_textbox(left=7500*min(lst1), top=6000*(sum(lst2) / len(lst2)), width=w, height=h)
                tf = txBox.text_frame
                tf.text = mytext
                #time.sleep(1)
                mytext=""
                lst1.clear()
                lst2.clear()
            mytext+=j[11]+" "
            m=x
            lst1.append(x)
            lst2.append(y)
        i+=1
    pic = slide.shapes.add_picture(img_path, left = left1, top =top, height=height1,width=width1)
    prs.save('output/hellohi.pptx')

if __name__ == "__main__":

    st.markdown(html_temp, unsafe_allow_html=True)
    button1= st.button("Import Video from PC ")
    #sentence = st.text_area('youtube link :', height=30)
    #button2= st.button("browse")

    def Browse():
        download_Directory = filedialog.askdirectory(initialdir="YOUR DIRECTORY PATH", title="Save Video")
        return download_Directory
        #download_Path.set(download_Directory)

    def Download(sentence,brow):
        #For downloading video from youtube
        '''link_of_the_video = sentence.get()
        download_Folder = brow.get()
        ydl_opts = {
            'format': 'best',
            'outtmpl':download_Folder + '/%(title)s.%(ext)s',
        }
        yt = link_of_the_video.strip()
        with youtube_dl.YoutubeDL(ydl_opts) as ydl:
            ydl.download([yt])
        messagebox.showinfo("SUCCESSFULLY","DOWNLOADED AND SAVED IN\n"+ download_Folder) '''
    
    def openFile():
        filepath = filedialog.askopenfilename()
        return filepath

    if button1 :
        filepath = filedialog.askopenfilename()
        print('video_path', filepath)
        output_folder_screenshot_path = initialize_output_folder(filepath)
        detect_unique_screenshots(filepath, output_folder_screenshot_path)


    #if button2 and sentence:
        '''brow=Browse()
        Download(sentence,brow)
        
        print('video_path', brow)
        output_folder_screenshot_path = initialize_output_folder(brow)
        detect_unique_screenshots(brow, output_folder_screenshot_path)
'''
